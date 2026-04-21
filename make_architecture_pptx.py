"""
Generates architecture.pptx — a single-slide architecture diagram
White & grey colour palette, no images needed.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Colours ───────────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF0, 0xF0, 0xF0)   # box fills
MID_GREY   = RGBColor(0xC0, 0xC0, 0xC0)   # box borders / arrows
DARK_GREY  = RGBColor(0x40, 0x40, 0x40)   # body text
CHARCOAL   = RGBColor(0x20, 0x20, 0x20)   # headings
ACCENT     = RGBColor(0x70, 0x70, 0x70)   # phase banners

# ── Slide setup ────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)

slide  = prs.slides.add_slide(prs.slide_layouts[6])   # blank
shapes = slide.shapes

# white background
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = WHITE


# ── Helper functions ───────────────────────────────────────────────────────────

def box(l, t, w, h,
        fill_rgb=LIGHT_GREY, border_rgb=MID_GREY,
        border_pt=1.0, radius=0.12):
    """Add a rounded rectangle."""
    shape = shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,   # MSO_SHAPE.ROUNDED_RECTANGLE = 5
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    # rounded rectangle magic number = 5
    shape.shape_type   # just access to be safe
    sp = shape._element
    spPr = sp.find(nsmap.qn('p:spPr'))
    prstGeom = spPr.find(nsmap.qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')
        avLst = prstGeom.find(nsmap.qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, nsmap.qn('a:avLst'))
        else:
            avLst.clear()
        gd = etree.SubElement(avLst, nsmap.qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', 'val 20000')   # roundness

    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = border_rgb
    shape.line.width = Pt(border_pt)
    shape.shadow.inherit = False
    return shape


def rect(l, t, w, h, fill_rgb=LIGHT_GREY, border_rgb=MID_GREY, border_pt=1.0):
    """Plain rectangle (for phase banners)."""
    shape = shapes.add_shape(
        1,  # AUTO_SHAPE → default is rectangle prstGeom=rect
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    sp = shape._element
    spPr = sp.find(nsmap.qn('p:spPr'))
    prstGeom = spPr.find(nsmap.qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'rect')

    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = border_rgb
    shape.line.width = Pt(border_pt)
    shape.shadow.inherit = False
    return shape


def label(shape, title, body="", title_sz=10, body_sz=8,
          title_rgb=CHARCOAL, body_rgb=DARK_GREY, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    p0 = tf.paragraphs[0]
    p0.alignment = align
    run = p0.add_run()
    run.text = title
    run.font.bold  = True
    run.font.size  = Pt(title_sz)
    run.font.color.rgb = title_rgb

    if body:
        for line in body.split('\n'):
            para = tf.add_paragraph()
            para.alignment = align
            r = para.add_run()
            r.text = line
            r.font.size = Pt(body_sz)
            r.font.color.rgb = body_rgb
            r.font.bold = False


def arrow(x1, y1, x2, y2):
    """Draw a connector arrow from (x1,y1) to (x2,y2) in inches."""
    from pptx.util import Inches
    from pptx.oxml.ns import qn
    from lxml import etree

    # Use a straight connector
    cxnSp = etree.SubElement(
        slide.shapes._spTree, qn('p:cxnSp')
    )
    nvCxnSpPr = etree.SubElement(cxnSp, qn('p:nvCxnSpPr'))
    cNvPr = etree.SubElement(nvCxnSpPr, qn('p:cNvPr'))
    cNvPr.set('id', str(len(shapes) + 200))
    cNvPr.set('name', f'conn{len(shapes)}')
    etree.SubElement(nvCxnSpPr, qn('p:cNvCxnSpPr'))
    etree.SubElement(nvCxnSpPr, qn('p:nvPr'))

    spPr = etree.SubElement(cxnSp, qn('p:spPr'))
    xfrm = etree.SubElement(spPr, qn('a:xfrm'))

    lx = min(x1, x2)
    ly = min(y1, y2)
    w  = abs(x2 - x1)
    h  = abs(y2 - y1)

    off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', str(int(Inches(lx))))
    off.set('y', str(int(Inches(ly))))
    ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', str(max(int(Inches(w)), 1)))
    ext.set('cy', str(max(int(Inches(h)), 1)))

    if x1 > x2 or y1 > y2:
        xfrm.set('flipH', '1') if x1 > x2 else None
        xfrm.set('flipV', '1') if y1 > y2 else None

    prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
    prstGeom.set('prst', 'straightConnector1')
    etree.SubElement(prstGeom, qn('a:avLst'))

    ln = etree.SubElement(spPr, qn('a:ln'))
    ln.set('w', str(int(Pt(1.2))))
    solidFill = etree.SubElement(ln, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', '909090')
    # arrowhead
    headEnd = etree.SubElement(ln, qn('a:headEnd'))
    tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'arrow')
    tailEnd.set('w', 'med')
    tailEnd.set('len', 'med')


# ═══════════════════════════════════════════════════════════════════════════════
# LAYOUT  (all in inches, slide = 16 × 9)
# ═══════════════════════════════════════════════════════════════════════════════

# ── Title ─────────────────────────────────────────────────────────────────────
title_shape = shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(15.4), Inches(0.5))
tf = title_shape.text_frame
p  = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Multi-Role Document Processing Pipeline Architecture"
r.font.size  = Pt(15)
r.font.bold  = True
r.font.color.rgb = CHARCOAL

# ── USER INPUT (top) ───────────────────────────────────────────────────────────
USER_L, USER_T, USER_W, USER_H = 5.5, 0.65, 5.0, 0.55
s = box(USER_L, USER_T, USER_W, USER_H, fill_rgb=RGBColor(0xE8,0xE8,0xE8))
label(s, "User uploads PDF + selects doc type",
      "e.g.  Banking → ISDA  |  Insurance → Invoice  |  Compliance → Report",
      title_sz=9, body_sz=7.5)

# OCR box
OCR_L, OCR_T, OCR_W, OCR_H = 6.5, 1.4, 3.0, 0.5
s = box(OCR_L, OCR_T, OCR_W, OCR_H, fill_rgb=RGBColor(0xE8,0xE8,0xE8))
label(s, "OCR  (cached by PDF hash)", title_sz=9)

# arrows user → OCR
arrow(8.0, 1.2, 8.0, 1.4)

# ── PHASE 1 — DOCUMENT MAP ────────────────────────────────────────────────────
P1_L, P1_T, P1_W, P1_H = 3.8, 2.1, 8.4, 2.1
s = rect(P1_L, P1_T, P1_W, P1_H, fill_rgb=RGBColor(0xF5,0xF5,0xF5),
         border_rgb=RGBColor(0xA0,0xA0,0xA0), border_pt=1.5)
label(s,
      "PHASE 1 — DOCUMENT MAP",
      "1 LLM call  ·  runs once per document\n\n"
      "Reads page summaries  →  returns:\n"
      '{ "parties_dates": {pages:[1,2], type:"paragraph"}  '
      '"rates": {pages:[5,6], type:"table"}\n'
      '  "covenants": {pages:[8,9], type:"paragraph"}  '
      '"fees": {pages:[12], type:"table"}  '
      '"definitions": {pages:[3,4], type:"paragraph"} }\n\n'
      "Tells you:  WHERE each topic is  +  TABLE or PARAGRAPH",
      title_sz=10, body_sz=7.5, align=PP_ALIGN.LEFT)

arrow(8.0, 1.9, 8.0, 2.1)

# ── PARALLEL ZONE label ───────────────────────────────────────────────────────
zone_label = shapes.add_textbox(Inches(0.2), Inches(4.35), Inches(2.0), Inches(0.3))
tf = zone_label.text_frame
p  = tf.paragraphs[0]
r  = p.add_run()
r.text = "Per key — 6 parallel workers"
r.font.size = Pt(7.5)
r.font.italic = True
r.font.color.rgb = DARK_GREY

arrow(8.0, 4.2, 8.0, 4.35)

# ── TABLE PATH ────────────────────────────────────────────────────────────────

# Tables Agent
TA_L, TA_T, TA_W, TA_H = 1.0, 4.5, 3.0, 1.5
s = box(TA_L, TA_T, TA_W, TA_H)
label(s, "TABLES AGENT",
      "Extracts raw table data\nH + V cells, structure\nmerged cells, cross-page",
      title_sz=9, body_sz=7.5)

# Rules Sub-Agent
RS_L, RS_T, RS_W, RS_H = 0.5, 6.3, 3.8, 1.45
s = box(RS_L, RS_T, RS_W, RS_H, fill_rgb=RGBColor(0xEA,0xEA,0xEA),
        border_rgb=RGBColor(0x88,0x88,0x88))
label(s, "RULES SUB-AGENT",
      "Input: table output + paragraphs surrounding the table\n"
      "Extracts: exact rules · conditions · thresholds\n"
      "          exceptions · context",
      title_sz=9, body_sz=7.5, align=PP_ALIGN.LEFT)

arrow(2.5, 6.0, 2.5, 6.3)

# ── PARAGRAPH PATH ────────────────────────────────────────────────────────────

# Doc-Type Sub-Agent
DT_L, DT_T, DT_W, DT_H = 5.2, 4.5, 3.4, 1.5
s = box(DT_L, DT_T, DT_W, DT_H)
label(s, "DOC-TYPE SUB-AGENT",
      "Prompt tailored to doc type:\n· ISDA     · Loan\n· Invoice  · Report",
      title_sz=9, body_sz=7.5)

# Terms Agent
TE_L, TE_T, TE_W, TE_H = 5.2, 6.3, 3.4, 1.45
s = box(TE_L, TE_T, TE_W, TE_H, fill_rgb=RGBColor(0xEA,0xEA,0xEA),
        border_rgb=RGBColor(0x88,0x88,0x88))
label(s, "TERMS AGENT",
      "Called when key demands rule/term\nextraction from prose\n\n"
      "Extracts: obligations · conditions\n"
      "          fixed terms · from clauses",
      title_sz=9, body_sz=7.5, align=PP_ALIGN.LEFT)

arrow(6.9, 6.0, 6.9, 6.3)

# ── DESCRIPTION AGENT ─────────────────────────────────────────────────────────
DA_L, DA_T, DA_W, DA_H = 10.3, 4.5, 3.2, 3.25
s = box(DA_L, DA_T, DA_W, DA_H, fill_rgb=RGBColor(0xEE,0xEE,0xEE),
        border_rgb=RGBColor(0x88,0x88,0x88))
label(s, "DESCRIPTION AGENT",
      "Always runs in parallel\n\n"
      "Searches definitions section\n\n"
      "Found  →  returns definition text\n"
      "Not found  →  \"\"  (never generates)",
      title_sz=9, body_sz=7.5)

# "type=table" / "type=paragraph" labels above boxes
for txt, lx, lt in [
    ("type = \"table\"",     1.6, 4.25),
    ("type = \"paragraph\"", 6.2, 4.25),
]:
    tb = shapes.add_textbox(Inches(lx), Inches(lt), Inches(2.5), Inches(0.22))
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r  = p.add_run()
    r.text = txt
    r.font.size = Pt(7.5)
    r.font.italic = True
    r.font.color.rgb = DARK_GREY

# ── VALIDATOR ─────────────────────────────────────────────────────────────────
VA_L, VA_T, VA_W, VA_H = 5.2, 7.95, 5.6, 0.78
s = rect(VA_L, VA_T, VA_W, VA_H, fill_rgb=RGBColor(0xE0,0xE0,0xE0),
         border_rgb=RGBColor(0x80,0x80,0x80), border_pt=1.5)
label(s, "PHASE 3 — VALIDATOR",
      "1. Format rule check (instant, no LLM)    "
      "2. LLM only if: value is null  or  format is wrong",
      title_sz=9.5, body_sz=7.5)

# ── OUTPUT ────────────────────────────────────────────────────────────────────
OUT_L, OUT_T, OUT_W, OUT_H = 4.8, 8.88, 6.4, 0.52
s = shapes.add_textbox(Inches(OUT_L), Inches(OUT_T), Inches(OUT_W), Inches(OUT_H))
tf = s.text_frame
p  = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r  = p.add_run()
r.text = "{ value,  format,  score,  page,  description,  rule_context,  found_in: table / paragraph / rule }"
r.font.size  = Pt(8)
r.font.color.rgb = DARK_GREY
r.font.italic = True

# ── ARROWS (phase 1 → parallel zone top; then down to validator) ──────────────
# Table agent ← phase 1
arrow(8.0, 4.2, 2.5, 4.5)
# Para agent ← phase 1
arrow(8.0, 4.2, 6.9, 4.5)
# Description ← phase 1
arrow(8.0, 4.2, 11.9, 4.5)

# Rules → Validator
arrow(2.5, 7.75, 8.0, 7.95)
# Terms → Validator
arrow(6.9, 7.75, 8.0, 7.95)
# Description → Validator
arrow(11.9, 7.75, 8.0, 7.95)

# Validator → Output
arrow(8.0, 8.73, 8.0, 8.88)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/home/mariem/deepdoctection_project/architecture.pptx"
prs.save(out)
print(f"Saved: {out}")
